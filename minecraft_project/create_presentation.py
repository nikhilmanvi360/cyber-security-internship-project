"""
PowerPoint Presentation Generator for Cybersecurity Internship Report
Creates a professional presentation from the keylogger project documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    DARK_BLUE = RGBColor(31, 78, 121)
    LIGHT_BLUE = RGBColor(68, 114, 196)
    ORANGE = RGBColor(237, 125, 49)
    GRAY = RGBColor(89, 89, 89)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Advanced Multi-Platform Keylogger"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Cybersecurity Research & Penetration Testing"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "Educational Project | November 2025\nPlatforms: Windows | Linux | Android"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.font.color.rgb = GRAY
    author_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Sophisticated cross-platform keylogger demonstrating:"
    
    points = [
        "Low-level system hooks for keystroke capture",
        "Military-grade AES-256 encryption",
        "Biometric authentication (facial recognition)",
        "Covert network exfiltration via Discord webhooks",
        "Multi-platform deployment (Windows, Linux, Android)",
        "Advanced persistence mechanisms"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 3: Architecture Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    components = [
        ("Main Controller", "Orchestrates all components, manages lifecycle"),
        ("KeyLogger Engine", "Captures keystrokes using pynput library"),
        ("Input Processor", "Processes and formats keystroke data"),
        ("Crypto Module", "AES-256 encryption/decryption"),
        ("Face Auth", "Biometric authentication via OpenCV"),
        ("Network Sender", "Exfiltrates data to Discord webhook")
    ]
    
    for comp, desc in components:
        p = tf.add_paragraph()
        p.text = f"{comp}: {desc}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 4: Code Statistics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Code Statistics"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    stats = [
        ("Total Lines of Code", "593 lines"),
        ("Programming Language", "Python 3"),
        ("Core Components", "7 modules"),
        ("External Libraries", "6 (pynput, cryptography, opencv, etc.)"),
        ("Platforms Supported", "3 (Windows, Linux, Android)"),
        ("Encryption Strength", "AES-256-CBC with HMAC-SHA256"),
        ("Executable Size", "12-25 MB (platform dependent)"),
        ("Memory Footprint", "~20-30 MB")
    ]
    
    for stat, value in stats:
        p = tf.add_paragraph()
        p.text = f"{stat}: {value}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 5: KeyLogger Engine Logic
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "KeyLogger Engine (logger.py)"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Core Functionality:"
    
    points = [
        "Event-driven architecture using pynput library",
        "Platform-specific hooks (Win32 API, X11/Xorg)",
        "Asynchronous key capture in separate thread",
        "Callback pattern for key event handling",
        "Exception handling prevents crashes",
        "Non-blocking design for continuous operation"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 6: Input Processor Logic
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Input Processor (processor.py)"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Processing Logic:"
    
    points = [
        "Classifies keys: KeyCode (regular chars) vs Key (special)",
        "Maintains list-based buffer for O(1) append operations",
        "Handles special keys: Space, Enter, Backspace, etc.",
        "Preserves ALL keystrokes including modifiers",
        "Reconstructs full text on each keystroke",
        "Provides statistics (character count, word count)"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 7: Encryption Module
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Encryption Module (crypto_utils.py)"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Cryptographic Security:"
    
    points = [
        "Fernet encryption (AES-256-CBC mode)",
        "256-bit symmetric key (military-grade)",
        "HMAC-SHA256 for authentication & integrity",
        "Base64 encoding for storage",
        "Timestamp inclusion prevents replay attacks",
        "Key stored in secret.key file"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 8: Face Authentication
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Biometric Authentication (face_auth.py)"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Face Recognition System:"
    
    points = [
        "Haar Cascade Classifier for face detection",
        "LBPH (Local Binary Patterns Histograms) recognition",
        "Captures 30 face samples for training",
        "Confidence threshold < 50 for authentication",
        "10-second timeout for authentication attempt",
        "Model stored in face_model.yml (4.3 MB)"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 9: Network Exfiltration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Network Exfiltration (network_sender.py)"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Data Exfiltration Techniques:"
    
    points = [
        "Discord webhook as C2 (Command & Control)",
        "HTTPS traffic blends with legitimate Discord usage",
        "Automated sending every 30 seconds",
        "Connectivity check before transmission",
        "Truncates logs to 1900 chars (Discord limit)",
        "User-Agent spoofing (mimics browser traffic)"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 10: Main Controller
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Main Controller (main.py)"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Orchestration Logic:"
    
    points = [
        "Detects portable execution (PyInstaller)",
        "Loads encryption key from secret.key",
        "Writes PID for process management",
        "Loads existing encrypted logs (persistence)",
        "Dual storage: encrypted (.enc) + plaintext (.txt)",
        "Continuous main loop with periodic exfiltration"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 11: Platform-Specific Implementations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Platform-Specific Implementations"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    platforms = [
        ("Windows", "PyInstaller, Win32 hooks, Registry autostart"),
        ("Linux", "PyInstaller, Xorg hooks, systemd services"),
        ("Android", "Buildozer, Accessibility Service, autostart on boot")
    ]
    
    for platform, details in platforms:
        p = tf.add_paragraph()
        p.text = f"{platform}:"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p2 = tf.add_paragraph()
        p2.text = details
        p2.level = 1
        p2.font.size = Pt(16)
    
    # Slide 12: Stealth Mechanisms
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Stealth & Anti-Detection"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Evasion Techniques:"
    
    points = [
        "No GUI (--noconsole flag, runs silently)",
        "Hidden installation directories (.minecraft_updater)",
        "Legitimate-sounding name ('Minecraft Launcher')",
        "HTTPS encryption for network traffic",
        "Minimal CPU usage (<1%)",
        "Small memory footprint (~20-30 MB)",
        "Blends with legitimate Discord traffic"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 13: Persistence Mechanisms
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Persistence Mechanisms"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    persistence = [
        ("Windows", "Registry Run key (HKCU\\...\\Run)"),
        ("Linux", "systemd user service with auto-restart"),
        ("Android", "BOOT_COMPLETED broadcast receiver"),
        ("All Platforms", "Encrypted log persistence across reboots")
    ]
    
    for platform, method in persistence:
        p = tf.add_paragraph()
        p.text = f"{platform}: {method}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 14: Security Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security Features Summary"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    features = [
        ("Encryption", "AES-256-CBC (military-grade)"),
        ("Authentication", "HMAC-SHA256 integrity verification"),
        ("Biometrics", "Facial recognition (LBPH algorithm)"),
        ("Stealth", "No GUI, hidden directories, process hiding"),
        ("Persistence", "Survives reboots, auto-restart on crash"),
        ("Exfiltration", "Covert HTTPS channel via Discord")
    ]
    
    for feature, detail in features:
        p = tf.add_paragraph()
        p.text = f"{feature}: {detail}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 15: Attack Vectors
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Attack Vectors & Deployment"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Social Engineering Scenarios:"
    
    points = [
        "Fake game launcher (targets gamers)",
        "Software bundle (bundled with legitimate software)",
        "USB drop attack (physical access)",
        "Trojanized installers (modified legitimate software)",
        "Phishing emails with malicious attachments"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 16: Detection Methods
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Detection & Defense"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Detection Techniques:"
    
    points = [
        "Monitor keyboard hook installations",
        "Detect registry autostart modifications",
        "Network traffic analysis (Discord webhooks)",
        "File integrity monitoring (%APPDATA% changes)",
        "Behavioral analysis (process patterns)",
        "Endpoint Detection & Response (EDR) tools"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 17: Technical Challenges
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Challenges & Solutions"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    challenges = [
        ("Cross-platform compatibility", "pynput library abstracts OS differences"),
        ("Stealth & anti-detection", "No GUI, hidden directories, HTTPS traffic"),
        ("Data persistence", "Encrypted logs loaded on startup"),
        ("Network reliability", "Connectivity checks, graceful failure"),
        ("Android restrictions", "Accessibility Service, battery optimization")
    ]
    
    for challenge, solution in challenges:
        p = tf.add_paragraph()
        p.text = f"{challenge}:"
        p.level = 0
        p.font.size = Pt(15)
        p.font.bold = True
        
        p2 = tf.add_paragraph()
        p2.text = solution
        p2.level = 1
        p2.font.size = Pt(14)
    
    # Slide 18: Ethical Considerations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Ethical & Legal Considerations"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Legal Framework:"
    
    points = [
        "Computer Fraud & Abuse Act (USA) - Up to 10 years",
        "GDPR violations (EU) - Fines up to €20 million",
        "Computer Misuse Act 1990 (UK) - Up to 2 years",
        "",
        "Legitimate Use Cases:",
        "• Parental monitoring (with consent)",
        "• Corporate security (with disclosure)",
        "• Penetration testing (with authorization)",
        "• Research & education (isolated environment)"
    ]
    
    for point in points:
        if point:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0 if point.startswith("•") else 1
            p.font.size = Pt(15)
    
    # Slide 19: Key Learnings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Learnings & Skills Acquired"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    learnings = [
        "Low-level system programming (keyboard hooks)",
        "Cryptography (AES-256, key management)",
        "Computer vision (face detection & recognition)",
        "Network programming (HTTP, webhooks)",
        "Cross-platform development (Windows/Linux/Android)",
        "Stealth techniques (process hiding, persistence)",
        "Build automation (PyInstaller, Buildozer)"
    ]
    
    for learning in learnings:
        p = tf.add_paragraph()
        p.text = learning
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 20: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    conclusion_text = """This project demonstrates comprehensive understanding of:

• Offensive security techniques
• Defensive cybersecurity measures
• Cross-platform software development
• Cryptographic implementations
• Network security & data exfiltration

Educational Purpose Only
Unauthorized deployment is illegal and unethical.
Always obtain proper authorization."""
    
    tf.text = conclusion_text
    tf.paragraphs[0].font.size = Pt(18)
    
    # Slide 21: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You\n\nQuestions?"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(54)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = DARK_BLUE
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # Add disclaimer
    disclaimer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    disclaimer_frame = disclaimer_box.text_frame
    disclaimer_frame.text = "Educational Research Project | November 2025"
    disclaimer_para = disclaimer_frame.paragraphs[0]
    disclaimer_para.font.size = Pt(16)
    disclaimer_para.font.color.rgb = GRAY
    disclaimer_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Cybersecurity_Keylogger_Presentation.pptx')
    print("✅ Presentation created successfully!")
    print("📄 File: Cybersecurity_Keylogger_Presentation.pptx")
    print(f"📊 Total Slides: {len(prs.slides)}")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not installed")
        print("Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
